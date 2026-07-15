#!/usr/bin/env python3
"""
Ingest course materials into vector database.

Usage:
    python scripts/ingestion/ingest_courses.py --course-dir data/raw/ASD-2024
    python scripts/ingestion/ingest_courses.py --all
"""

import argparse
import json
import sys
from pathlib import Path
from typing import Dict, List, Optional

from langchain_text_splitters import RecursiveCharacterTextSplitter
from tqdm import tqdm

# Add project root to path
project_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(project_root))

from src.core.config import settings
from src.core.config.logging import LoggerAdapter
from src.core.embeddings import get_embedding_service
from src.core.rag.bm25_encoder import BM25Encoder
from src.data_pipeline.storage import get_qdrant_storage

logger = LoggerAdapter(__name__)


def load_document(file_path: Path) -> str:
    """
    Load document from file.

    Args:
        file_path: Path to document

    Returns:
        str: Document text
    """
    suffix = file_path.suffix.lower()

    try:
        if suffix == ".txt":
            return file_path.read_text(encoding="utf-8", errors="ignore")

        elif suffix == ".pdf":
            try:
                from pypdf import PdfReader
            except ImportError:
                from PyPDF2 import PdfReader

            reader = PdfReader(file_path)
            pages = []
            for page in reader.pages:
                t = page.extract_text()
                if t and t.strip():
                    pages.append(t.strip())
            return "\n\n".join(pages)

        elif suffix == ".docx":
            from docx import Document

            doc = Document(file_path)
            text = "\n\n".join(para.text for para in doc.paragraphs)
            return text

        elif suffix == ".pptx":
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n\n".join(text_parts)

        else:
            logger.warning("Unsupported file type", file=str(file_path), suffix=suffix)
            return ""

    except Exception as e:
        logger.error("Failed to load document", file=str(file_path), error=str(e))
        return ""


def load_enrichment_cache(cache_path: str) -> Dict[str, str]:
    """Load document summaries from contextual enrichment cache."""
    path = Path(cache_path)
    if not path.exists():
        return {}
    cache = {}
    with open(path, encoding="utf-8") as f:
        for line in f:
            if line.strip():
                rec = json.loads(line)
                cache[rec["file"]] = rec["summary"]
    return cache


def ingest_course(
    course_dir: Path,
    enrichment_cache: Optional[Dict[str, str]] = None,
    bm25_collector: Optional[List] = None,
    hybrid_mode: bool = False,
) -> None:
    """
    Ingest all documents from course directory.

    Args:
        course_dir: Course directory path
        enrichment_cache: Optional {file_path: summary} for contextual enrichment
        bm25_collector: If provided, collect (chunk_text, indices, values) for BM25
        hybrid_mode: If True, use named dense vectors + sparse vectors
    """
    course_id = course_dir.name
    logger.info("Ingesting course", course_id=course_id, path=str(course_dir))

    supported_extensions = [".pdf", ".docx", ".pptx", ".txt"]
    files: List[Path] = []
    for ext in supported_extensions:
        files.extend(course_dir.rglob(f"*{ext}"))

    if not files:
        logger.warning("No documents found", course_dir=str(course_dir))
        return

    logger.info("Found documents", count=len(files))

    embedding_service = get_embedding_service()
    storage = get_qdrant_storage()

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.rag_chunk_size,
        chunk_overlap=settings.rag_chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    all_chunks_original = []   # stored in payload (for display)
    all_chunks_to_embed = []   # may include context prefix
    all_payloads = []

    for file_path in tqdm(files, desc=f"Processing {course_id}"):
        text = load_document(file_path)
        if not text.strip():
            continue

        chunks = text_splitter.split_text(text)
        summary = (enrichment_cache or {}).get(str(file_path), "") if enrichment_cache else ""

        for i, chunk in enumerate(chunks):
            payload = {
                "text": chunk,
                "course_id": course_id,
                "course_name": course_id.replace("-", " "),
                "file_path": file_path.name,
                "chunk_index": i,
                "topic": file_path.stem,
            }
            # Prepend document summary for richer embeddings (contextual enrichment)
            embed_text = f"{summary}\n\n{chunk}" if summary else chunk

            all_chunks_original.append(chunk)
            all_chunks_to_embed.append(embed_text)
            all_payloads.append(payload)

    if not all_chunks_to_embed:
        logger.warning("No chunks extracted", course_id=course_id)
        return

    logger.info("Extracted chunks", total_chunks=len(all_chunks_to_embed))
    enriched = bool(enrichment_cache) and any(str(f) in enrichment_cache for f in files)
    print(f"  Contextual enrichment: {'✅' if enriched else '—'}")

    print(f"Generating embeddings for {len(all_chunks_to_embed)} chunks...")
    embeddings = embedding_service.encode_documents(
        all_chunks_to_embed,
        batch_size=settings.batch_size_embedding,
        show_progress_bar=True,
    )

    if hybrid_mode and bm25_collector is not None:
        # Collect texts for BM25 fitting (will be done globally after all courses)
        for chunk, emb, payload in zip(all_chunks_to_embed, embeddings, all_payloads):
            bm25_collector.append((chunk, emb, payload))
    else:
        print("Uploading to Qdrant...")
        storage.upsert_points(
            vectors=embeddings,
            payloads=all_payloads,
            batch_size=settings.batch_size_upload,
        )

    logger.info("Course ingestion completed", course_id=course_id, chunks=len(all_chunks_to_embed))
    print(f"✅ Ingested {course_id}: {len(all_chunks_to_embed)} chunks")


def main() -> None:
    """Main function."""
    parser = argparse.ArgumentParser(description="Ingest course materials")
    parser.add_argument(
        "--course-dir",
        type=Path,
        help="Course directory to ingest",
    )
    parser.add_argument(
        "--all",
        action="store_true",
        help="Ingest all courses in data/raw/",
    )
    parser.add_argument(
        "--data-dir",
        type=Path,
        default=Path(settings.data_raw_path),
        help="Data directory (default: data/raw)",
    )
    parser.add_argument(
        "--wipe",
        action="store_true",
        help="Delete and recreate Qdrant collection before ingesting (use after model change)",
    )
    parser.add_argument(
        "--hybrid",
        action="store_true",
        help="Enable hybrid dense+sparse (BM25) search — requires --wipe on first use",
    )
    parser.add_argument(
        "--enrichment-cache",
        default="data/enrichment_cache.jsonl",
        help="Path to contextual enrichment cache (from contextual_enrichment.py)",
    )
    parser.add_argument(
        "--bm25-vocab",
        default="data/bm25_vocab.json",
        help="Where to save BM25 vocabulary (used at query time)",
    )

    args = parser.parse_args()

    enrichment_cache = load_enrichment_cache(args.enrichment_cache)
    if enrichment_cache:
        print(f"Loaded enrichment summaries for {len(enrichment_cache)} documents")
    else:
        print("No enrichment cache found — embedding raw chunks")

    if args.wipe:
        from src.data_pipeline.storage import get_qdrant_storage
        storage = get_qdrant_storage()
        print("Wiping Qdrant collection...")
        storage.delete_collection()
        if args.hybrid:
            storage.create_collection_hybrid()
            print("Hybrid collection (dense+sparse) recreated.")
        else:
            storage.create_collection()
            print("Collection recreated.")

    if args.all:
        # Ingest all courses
        data_dir = args.data_dir
        if not data_dir.exists():
            print(f"❌ Data directory not found: {data_dir}")
            sys.exit(1)

        # Find all course directories
        course_dirs = [d for d in data_dir.iterdir() if d.is_dir() and not d.name.startswith(".")]

        if not course_dirs:
            print(f"❌ No course directories found in {data_dir}")
            sys.exit(1)

        print(f"Found {len(course_dirs)} courses:")
        for course_dir in course_dirs:
            print(f"  - {course_dir.name}")
        print()

        bm25_collector = [] if args.hybrid else None

        for course_dir in course_dirs:
            try:
                ingest_course(
                    course_dir,
                    enrichment_cache=enrichment_cache or None,
                    bm25_collector=bm25_collector,
                    hybrid_mode=args.hybrid,
                )
            except Exception as e:
                logger.error("Failed to ingest course", course_id=course_dir.name, error=str(e))
                print(f"❌ Error ingesting {course_dir.name}: {e}")

        # Hybrid mode: fit BM25 on full corpus, compute sparse vectors, upload
        if args.hybrid and bm25_collector:
            from src.data_pipeline.storage import get_qdrant_storage
            storage = get_qdrant_storage()

            print(f"\nFitting BM25 on {len(bm25_collector)} chunks...")
            all_texts = [item[0] for item in bm25_collector]
            encoder = BM25Encoder()
            encoder.fit(all_texts)
            encoder.save(args.bm25_vocab)
            print(f"BM25 vocab size: {encoder.vocab_size} — saved to {args.bm25_vocab}")

            dense_vecs, sparse_idx, sparse_val, payloads = [], [], [], []
            for text, emb, payload in bm25_collector:
                indices, values = encoder.encode_document(text)
                dense_vecs.append(emb)
                sparse_idx.append(indices)
                sparse_val.append(values)
                payloads.append(payload)

            print(f"Uploading {len(dense_vecs)} hybrid points to Qdrant...")
            storage.upsert_points_hybrid(
                dense_vectors=dense_vecs,
                sparse_indices=sparse_idx,
                sparse_values=sparse_val,
                payloads=payloads,
                batch_size=100,
            )

        print("\n✅ All courses ingested!")

    elif args.course_dir:
        if not args.course_dir.exists():
            print(f"❌ Course directory not found: {args.course_dir}")
            sys.exit(1)

        try:
            ingest_course(
                args.course_dir,
                enrichment_cache=enrichment_cache or None,
            )
        except Exception as e:
            logger.error("Failed to ingest course", error=str(e), exc_info=True)
            print(f"❌ Error: {e}")
            sys.exit(1)

    else:
        parser.print_help()
        sys.exit(1)


if __name__ == "__main__":
    main()
