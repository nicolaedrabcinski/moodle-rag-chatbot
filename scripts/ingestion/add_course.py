#!/usr/bin/env python3
"""
Add a new course to the RAG system: ingest materials into Qdrant + create Moodle course.

Usage:
    # Add course materials (directory auto-created if missing):
    python scripts/ingestion/add_course.py --course-id ASD-2026 --materials-dir /path/to/pdfs

    # Just ingest already-placed materials from data/raw/ASD-2026/:
    python scripts/ingestion/add_course.py --course-id ASD-2026

    # Skip Moodle course creation:
    python scripts/ingestion/add_course.py --course-id ASD-2026 --no-moodle

    # List known courses:
    python scripts/ingestion/add_course.py --list
"""

import argparse
import json
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple

project_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(project_root))

from src.core.config import settings
from src.core.config.logging import LoggerAdapter
from src.core.embeddings import get_embedding_service
from src.core.rag.bm25_encoder import BM25Encoder
from src.data_pipeline.storage import get_qdrant_storage

logger = LoggerAdapter(__name__)

COURSES_CONFIG = project_root / "data" / "courses_config.json"
BM25_VOCAB_PATH = project_root / "data" / "bm25_vocab.json"
RAW_DATA_DIR = project_root / "data" / "raw"

SUPPORTED_EXTENSIONS = [".pdf", ".txt", ".docx", ".pptx"]


def load_courses_config() -> Dict:
    if COURSES_CONFIG.exists():
        return json.loads(COURSES_CONFIG.read_text())
    return {}


def load_document(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    try:
        if suffix == ".txt":
            return file_path.read_text(encoding="utf-8", errors="ignore")
        elif suffix == ".pdf":
            try:
                from pypdf import PdfReader
            except ImportError:
                from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            pages = [p.extract_text() for p in reader.pages]
            return "\n\n".join(t.strip() for t in pages if t and t.strip())
        elif suffix == ".docx":
            from docx import Document
            doc = Document(file_path)
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n\n".join(parts)
    except Exception as e:
        print(f"  ⚠ Could not read {file_path.name}: {e}")
    return ""


def ingest_course_materials(
    course_id: str,
    course_name: str,
    course_dir: Path,
) -> int:
    """Ingest course materials into Qdrant. Returns number of chunks ingested."""
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    files = []
    for ext in SUPPORTED_EXTENSIONS:
        files.extend(course_dir.rglob(f"*{ext}"))

    if not files:
        print(f"  ❌ No supported files found in {course_dir}")
        print(f"     Supported formats: {', '.join(SUPPORTED_EXTENSIONS)}")
        return 0

    print(f"  Found {len(files)} file(s): {[f.name for f in files]}")

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.rag_chunk_size,
        chunk_overlap=settings.rag_chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    all_texts: List[str] = []
    all_payloads: List[Dict] = []

    for file_path in files:
        print(f"  Processing {file_path.name}...", end=" ", flush=True)
        text = load_document(file_path)
        if not text.strip():
            print("empty, skipped")
            continue
        chunks = text_splitter.split_text(text)
        print(f"{len(chunks)} chunks")
        for i, chunk in enumerate(chunks):
            all_texts.append(chunk)
            all_payloads.append({
                "text": chunk,
                "course_id": course_id,
                "course_name": course_name,
                "file_path": file_path.name,
                "chunk_index": i,
                "topic": file_path.stem,
            })

    if not all_texts:
        print("  ❌ No text could be extracted")
        return 0

    print(f"\n  Generating embeddings for {len(all_texts)} chunks...")
    embedding_service = get_embedding_service()
    embeddings = embedding_service.encode_documents(
        all_texts,
        batch_size=settings.batch_size_embedding,
        show_progress_bar=True,
    )

    storage = get_qdrant_storage()

    # Try hybrid upload with existing BM25 vocab; fall back to dense-only
    if BM25_VOCAB_PATH.exists():
        print("  Loading existing BM25 vocabulary...")
        encoder = BM25Encoder()
        encoder.load(str(BM25_VOCAB_PATH))

        sparse_idx, sparse_val = [], []
        for text in all_texts:
            idx, val = encoder.encode_document(text)
            sparse_idx.append(idx)
            sparse_val.append(val)

        print("  Uploading hybrid (dense+sparse) points to Qdrant...")
        storage.upsert_points_hybrid(
            dense_vectors=embeddings,
            sparse_indices=sparse_idx,
            sparse_values=sparse_val,
            payloads=all_payloads,
            batch_size=100,
        )
        print("  ℹ  BM25 vocab NOT re-fitted. New tokens will be searchable via dense only.")
        print("     Run 'python scripts/ingestion/ingest_courses.py --all --hybrid' to re-fit.")
    else:
        print("  Uploading dense-only points to Qdrant...")
        storage.upsert_points(
            vectors=embeddings,
            payloads=all_payloads,
            batch_size=settings.batch_size_upload,
        )

    return len(all_texts)


def create_moodle_course(course_id: str, course_name: str, description: str) -> bool:
    """Create course in Moodle via admin CLI."""
    print(f"\n  Creating Moodle course '{course_id}'...")
    cmd = [
        "docker", "exec", "fcim-moodle",
        "php", "/var/www/html/admin/cli/create_course.php",
        f"--shortname={course_id}",
        f"--fullname={course_name}",
        f"--summary={description}",
        "--visible=1",
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
        if result.returncode == 0:
            print(f"  ✅ Moodle course created: {course_id}")
            return True
        else:
            # May already exist — that's fine
            stderr = result.stderr + result.stdout
            if "already exists" in stderr.lower() or "duplicate" in stderr.lower():
                print(f"  ℹ  Moodle course already exists: {course_id}")
                return True
            print(f"  ⚠ Moodle CLI returned non-zero: {stderr[:200]}")
            # Try DB-level fallback
            return create_moodle_course_db(course_id, course_name, description)
    except Exception as e:
        print(f"  ⚠ Could not create Moodle course via CLI: {e}")
        return create_moodle_course_db(course_id, course_name, description)


def create_moodle_course_db(course_id: str, course_name: str, description: str) -> bool:
    """Fallback: create course directly in MariaDB."""
    import time
    ts = int(time.time())
    sql = (
        f"INSERT IGNORE INTO mdl_course (category, fullname, shortname, summary, visible, timecreated, timemodified, format) "
        f"VALUES (1, '{course_name}', '{course_id}', '{description}', 1, {ts}, {ts}, 'topics');"
    )
    cmd = [
        "docker", "exec", "fcim-moodle-db",
        "mysql", "-umoodle", "-pMoodleDB@2026!", "moodle", "-e", sql,
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=15)
        if result.returncode == 0:
            print(f"  ✅ Moodle course created via DB: {course_id}")
            return True
        else:
            print(f"  ⚠ DB insert failed: {result.stderr[:100]}")
            return False
    except Exception as e:
        print(f"  ⚠ DB fallback failed: {e}")
        return False


def clear_redis_cache() -> None:
    """Flush Redis cache so stale responses are invalidated."""
    cmd = ["docker", "exec", "fcim-redis", "redis-cli", "FLUSHDB"]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=10)
        if result.returncode == 0:
            print("  ✅ Redis cache cleared")
        else:
            print(f"  ⚠ Redis flush failed: {result.stderr[:80]}")
    except Exception as e:
        print(f"  ⚠ Could not clear Redis: {e}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Add a course to the RAG chatbot")
    parser.add_argument("--course-id", help="Course ID, e.g. ASD-2026")
    parser.add_argument("--course-name", help="Human-readable name (auto-detected from config)")
    parser.add_argument(
        "--materials-dir",
        type=Path,
        help="Source directory with PDF/DOCX files to copy into data/raw/<course-id>/",
    )
    parser.add_argument(
        "--no-moodle",
        action="store_true",
        help="Skip Moodle course creation",
    )
    parser.add_argument(
        "--no-cache-clear",
        action="store_true",
        help="Skip Redis cache flush",
    )
    parser.add_argument(
        "--list",
        action="store_true",
        help="List configured courses and their ingest status",
    )
    args = parser.parse_args()

    config = load_courses_config()

    if args.list:
        print("Configured courses:")
        for cid, meta in config.items():
            raw_dir = RAW_DATA_DIR / cid
            status = "✅ data present" if raw_dir.exists() and any(raw_dir.iterdir()) else "❌ no data"
            print(f"  {cid:15s}  {meta['name']:<40s}  {status}")
        return

    if not args.course_id:
        parser.print_help()
        sys.exit(1)

    course_id = args.course_id
    meta = config.get(course_id, {})
    course_name = args.course_name or meta.get("name") or course_id.replace("-", " ")
    description = meta.get("description", "")

    print(f"\n{'='*60}")
    print(f"Adding course: {course_id}")
    print(f"Name:          {course_name}")
    print(f"{'='*60}\n")

    # Step 1: copy materials if --materials-dir given
    course_dir = RAW_DATA_DIR / course_id
    if args.materials_dir:
        if not args.materials_dir.exists():
            print(f"❌ Materials directory not found: {args.materials_dir}")
            sys.exit(1)
        course_dir.mkdir(parents=True, exist_ok=True)
        files_to_copy = [
            f for f in args.materials_dir.iterdir()
            if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS
        ]
        if not files_to_copy:
            print(f"❌ No supported files in {args.materials_dir}")
            sys.exit(1)
        print(f"Copying {len(files_to_copy)} file(s) to {course_dir}...")
        for f in files_to_copy:
            shutil.copy2(f, course_dir / f.name)
            print(f"  Copied: {f.name}")

    if not course_dir.exists() or not any(course_dir.iterdir()):
        print(f"❌ No materials found in {course_dir}")
        print(f"   Place PDF/DOCX/TXT files there and re-run, or use --materials-dir /path/to/files")
        sys.exit(1)

    # Step 2: ingest into Qdrant
    print(f"\n[1/3] Ingesting materials into Qdrant...")
    n_chunks = ingest_course_materials(course_id, course_name, course_dir)
    if n_chunks == 0:
        print("❌ Nothing ingested. Aborting.")
        sys.exit(1)
    print(f"  ✅ {n_chunks} chunks indexed")

    # Step 3: create Moodle course
    if not args.no_moodle:
        print(f"\n[2/3] Creating Moodle course...")
        create_moodle_course(course_id, course_name, description)
    else:
        print("\n[2/3] Skipping Moodle course creation (--no-moodle)")

    # Step 4: clear cache
    if not args.no_cache_clear:
        print(f"\n[3/3] Clearing Redis cache...")
        clear_redis_cache()
    else:
        print("\n[3/3] Skipping cache clear (--no-cache-clear)")

    print(f"\n{'='*60}")
    print(f"✅ Course '{course_id}' ({course_name}) added successfully!")
    print(f"   Chunks in Qdrant: {n_chunks}")
    print(f"\nTo use in chatbot, pass course_id='{course_id}' in the request,")
    print(f"or leave it empty to search across all courses.")
    print(f"{'='*60}\n")


if __name__ == "__main__":
    main()
